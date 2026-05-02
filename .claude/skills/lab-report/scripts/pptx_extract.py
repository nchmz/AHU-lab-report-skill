"""
从 PPT 文件中提取文本内容和图片资源。

依赖：pip install python-pptx

用法：
    python pptx_extract.py <pptx文件路径> [输出目录]

输出：
    - text.md：提取的所有文本内容（按幻灯片分节）
    - images/：提取的所有图片
"""

import sys
import os
from pathlib import Path

sys.stdout.reconfigure(encoding='utf-8')

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print("Error: python-pptx library is required")
    print("Run: pip install python-pptx")
    sys.exit(1)


def extract_text(prs: Presentation) -> str:
    """提取 PPT 中所有文本内容，按幻灯片分节返回 Markdown 格式。"""
    lines = []

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}\n")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        # 根据层级添加缩进
                        indent = "  " * paragraph.level
                        lines.append(f"{indent}{text}")

            if shape.has_table:
                table = shape.table
                lines.append("")
                # 表头
                header = " | ".join(
                    cell.text.strip() for cell in table.rows[0].cells
                )
                lines.append(f"| {header} |")
                lines.append("|" + " --- |" * len(table.columns))
                # 数据行
                for idx, row in enumerate(table.rows):
                    if idx == 0:
                        continue
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells
                    )
                    lines.append(f"| {row_text} |")
                lines.append("")

        lines.append("")

    return "\n".join(lines)


def extract_images(prs: Presentation, output_dir: Path) -> list[str]:
    """提取 PPT 中所有图片，返回提取的文件路径列表。"""
    images_dir = output_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    extracted = []
    image_count = 0

    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image = shape.image
                ext = image.content_type.split("/")[-1]
                # 处理常见的图片格式映射
                ext_map = {
                    "jpeg": "jpg",
                    "x-emf": "emf",
                    "x-wmf": "wmf",
                    "tiff": "tif",
                }
                ext = ext_map.get(ext, ext)

                image_count += 1
                filename = f"slide{i}_img{image_count}.{ext}"
                filepath = images_dir / filename

                with open(filepath, "wb") as f:
                    f.write(image.blob)

                extracted.append(str(filepath))
                print(f"  Extracted: {filename}")

    return extracted


def main():
    if len(sys.argv) < 2:
        print("Usage: python pptx_extract.py <pptx_file_path> [output_directory]")
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    if not pptx_path.exists():
        print(f"Error: file not found - {pptx_path}")
        sys.exit(1)

    output_dir = Path(sys.argv[2]) if len(sys.argv) > 2 else Path("pptx_output")
    output_dir.mkdir(parents=True, exist_ok=True)

    print(f"Processing: {pptx_path}")
    prs = Presentation(str(pptx_path))

    # 提取文本
    print("\n[1/2] Extracting text content...")
    text_content = extract_text(prs)
    text_file = output_dir / "text.md"
    with open(text_file, "w", encoding="utf-8") as f:
        f.write(text_content)
    print(f"  Saved to: {text_file}")

    # 提取图片
    print("\n[2/2] Extracting image resources...")
    images = extract_images(prs, output_dir)
    print(f"  Extracted {len(images)} images total")

    # 输出摘要
    print(f"\nDone! Output directory: {output_dir}")
    print(f"  - Text: {text_file}")
    print(f"  - Images: {output_dir / 'images/'}")


if __name__ == "__main__":
    main()
